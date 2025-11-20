from pathlib import Path
from pptx import Presentation
import re, json
from shutil import rmtree
from tqdm import tqdm


PPTX2MD_VERSION = "0.1a.post1"


def generate_pptxt(
    path: Path, output_dir: Path, img_dir: Path, obsidian_compat: bool = False
) -> None:
    shape_table = {
        "title": ([8.83], [5.25, 5.33]),
        "header": ([9.0], [0.92, 1.0]),
        "paragraph": ([9.0, 9.17, 5.17, 9.08, 8.5], [5.85, 5.83]),
    }

    filename = path.stem
    output_loc = output_dir.joinpath(Path(f"{filename}.pptxt"))
    with open(output_loc, "w", encoding="utf-8") as f:
        f.write("[START PRESENTATION]\n")
        slides = Presentation(path).slides
        for slide_num, slide in enumerate(slides):
            f.write(f"\t[START SLIDE {slide_num}]\n")
            for shape in slide.shapes:

                # handle text
                if shape.has_text_frame and "Rectangle" in shape.name:
                    dim = (round(shape.width.inches, 2), round(shape.height.inches, 2))
                    shape_type = None
                    for k, v in shape_table.items():
                        if dim[0] in v[0] and dim[1] in v[1]:
                            shape_type = k
                    if not shape_type:
                        shape_type = "other"
                    highlight = "\x1b[7m" if shape_type == "other" else "\x1b[2m"

                    # print(f"{highlight}Name: {shape.name}\nDimensions: {dim}\nType: {shape_type}\nSlide: {slide_num+1} of {len(slides)}\nPath: {path}\n\x1b[0m")

                    f.write("\t\t")
                    match shape_type:
                        case "title":
                            f.write("[START TITLE]")
                        case "header":
                            f.write("[START HEADER]")
                        case "paragraph":
                            f.write("[START PARAGRAPH]")
                        case "other":
                            f.write("[START OTHER]")
                    f.write("\n")

                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            f.write(f"\t\t\t{run.text}\n")

                    f.write("\t\t")
                    match shape_type:
                        case "title":
                            f.write("[END TITLE]")
                        case "header":
                            f.write("[END HEADER]")
                        case "paragraph":
                            f.write("[END PARAGRAPH]")
                        case "other":
                            f.write("[END OTHER]")
                    f.write("\n")

                # handle images
                if shape.shape_type.name == "PICTURE":
                    f.write("\t\t[IMAGE ")
                    image = shape.image
                    # img_dir = output_dir.joinpath("images")
                    if not img_dir.exists():
                        img_dir.mkdir()

                    imgfn = image.filename
                    if imgfn.startswith("image."):
                        imgfn = f"{image.sha1}.{image.ext}"
                    img_path = img_dir.joinpath(imgfn).absolute()
                    obsidian_safe_path = "/".join(img_path.parts[-2:])
                    with open(img_path, "wb") as img:
                        # with open(imgfn, "wb") as img:
                        img.write(image.blob)
                    if not obsidian_compat:
                        f.write(f"{img_path}]\n")
                    else:
                        f.write(f"{obsidian_safe_path}]\n")

                # handle tables
                if shape.shape_type.name == "TABLE":
                    table = shape.table
                    headers = []
                    rows = []
                    for row_num, row in enumerate(table.rows):
                        _row = []
                        for cell_num, cell in enumerate(row.cells):
                            if row_num == 0:
                                headers.append(cell.text)
                            else:
                                _row.append(cell.text)
                        rows.append(_row)
                    table_obj = json.dumps({"headers": headers, "rows": rows})

                    f.write(f"\t\t[TABLE {table_obj}]\n")

            f.write(f"\t[END SLIDE {slide_num}]\n")
        f.write("[END PRESENTATION]")


def parse_pptxt(path: Path) -> dict:
    presentation = {"presentation": {"slides": []}}

    with open(path, "r", encoding="utf-8") as f:
        txt = f.read()

    slide_pat = re.compile(r"\[START SLIDE \d+\]\s+(.+?)\s+\[END SLIDE \d+\]", re.S)
    segment_pat = re.compile(
        r"(?:\[(IMAGE|TABLE) (.+?)(?<![\"\[\]])\])|(?:\[START (TITLE|HEADER|PARAGRAPH|IMAGE|OTHER)\]\s+(.+?)\s+\[END (?:TITLE|HEADER|PARAGRAPH|IMAGE|OTHER)\])",
        re.S,
    )

    slides = slide_pat.findall(txt)
    for slide in slides:
        slide_obj = {"segments": []}
        segments = segment_pat.findall(slide)
        # print(slide+"\n-----------")
        for segment in segments:
            segtype = segment[0] if segment[0] else segment[2]
            content = segment[1] if segment[1] else segment[3]
            content = content.replace("\t", "")
            segment_obj = {"type": segtype.lower(), "content": content}
            # print(segment_obj)
            slide_obj["segments"].append(segment_obj)
        presentation["presentation"]["slides"].append(slide_obj)

    return presentation


def parse_json_table(table: str) -> str:
    # print("\n"+table+"\n")
    table = json.loads(table)
    headers = table.get("headers")
    rows = table.get("rows")
    s = "\n|"
    for header in headers:
        s += f"{header}|"
    s += "\n|"
    for _ in range(len(headers)):
        s += ":---:|"
    for row in rows:
        for cell in row:
            s += f"{cell}|"
        s += "\n|"
    return s[:-2]


def pptxt_to_md(path: Path, output_dir: Path) -> None:
    with open(path, "r", encoding="utf-8") as f:
        presentation = json.load(f)
    presentation = presentation["presentation"]
    slides = presentation["slides"]
    pres = ""
    for slide_num, slide in enumerate(slides):
        is_new_topic = True
        segments = slide.get("segments")
        s = ""
        for segment in segments:
            seg = ""
            content = segment.get("content")
            segtype = segment.get("type")
            if content:
                match segtype:
                    case "title":
                        seg += "# "
                    case "header":
                        prev_headers = [
                            prev_seg.get("content")
                            for prev_seg in slides[slide_num - 1].get("segments")
                            if prev_seg.get("type") == "header"
                        ]
                        is_new_topic = not any([h == content for h in prev_headers])
                        if is_new_topic:
                            seg += "## "
                    case "paragraph":
                        seg += ""
                    case "image":
                        seg += r'<img src="'
                        # seg += "![["
                    case "table":
                        seg += parse_json_table(content)
                    case "other":
                        seg += ""

            if segtype == "header":
                # prev_headers = [prev_seg.get("content") for prev_seg in slides[slide_num-1].get("segments") if prev_seg.get("type") == "header"]
                # is_new_topic = not any([h == content for h in prev_headers])
                if is_new_topic:
                    seg += content
            elif segtype == "image":
                seg += content
                seg += r'">'
            elif segtype == "table":
                pass
            else:
                seg += content
            seg = seg.replace("\t", "")
            s += seg + "\n"
        pres += s
        if is_new_topic:
            pres += "\n---\n"
    output_loc = output_dir.joinpath(f"{path.stem}.md")
    with open(output_loc, "w", encoding="utf-8") as f:
        f.write(pres)
        f.write(f"<br><br><em><sub>generated by **pptx2md-re v{PPTX2MD_VERSION}**</sub></em>")


def process_pptx(
    slides_dir: Path,
    pptxt_dir: Path,
    json_dir: Path,
    md_dir: Path,
    img_dir: Path,
    keep_pptxt: bool = True,
    keep_json: bool = True,
    keep_md: bool = True,
    obsidian_compat: bool = False,
) -> None:

    for d in [pptxt_dir, json_dir, md_dir]:
        if d.exists():
            rmtree(d)
        d.mkdir()

    for path in tqdm(list(slides_dir.glob("*.pptx")), desc="parsing pptx -> pptxt"):
        generate_pptxt(path, pptxt_dir, img_dir, obsidian_compat)

    for path in tqdm(list(pptxt_dir.glob("*.pptxt")), desc="parsing pptxt -> json"):
        presentation = parse_pptxt(path)
        output_loc = json_dir.joinpath(f"{path.stem}.json")
        with open(output_loc, "w", encoding="utf-8") as f:
            json.dump(presentation, f, indent=4, ensure_ascii=True)

    for path in tqdm(list(json_dir.glob("*.json")), desc="parsing json -> md"):
        pptxt_to_md(path, md_dir)

    if not keep_pptxt:
        rmtree(pptxt_dir)
    if not keep_json:
        rmtree(json_dir)
    if not keep_md:
        rmtree(md_dir)


if __name__ == "__main__":
    slides_dir = Path(r"C:\Users\matcha\Documents\GRCC\microeconomics\Class Slides")
    pptxt_dir = Path(r"./pptxt")
    json_dir = Path(r"./json")
    md_dir = Path(r"C:\Users\matcha\Documents\Obsidian\MVault\pptxt")
    img_dir = Path(r"C:\Users\matcha\Documents\Obsidian\MVault\pptxt\images")
    # md_dir = Path(r".\md")
    # img_dir = Path(r".\md\images")
    obsidian_compat = True
    process_pptx(
        slides_dir=slides_dir,
        pptxt_dir=pptxt_dir,
        json_dir=json_dir,
        md_dir=md_dir,
        img_dir=img_dir,
        keep_json=True,
        keep_pptxt=True,
        obsidian_compat=obsidian_compat
    )
    # main()
